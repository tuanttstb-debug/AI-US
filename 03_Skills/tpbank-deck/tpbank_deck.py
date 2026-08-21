# -*- coding: utf-8 -*-
"""
tpbank_deck — thư viện dựng PowerPoint theo chuẩn nhận diện TPBank.

Trích từ template gốc "Báo cáo dự án Biz Anchor Lending". Dùng lại asset brand
thật (logo, dải accent, nền tím) trong ./assets. Khổ 16:9 (33.87 x 19.05 cm).

Cách dùng nhanh:
    from tpbank_deck import Deck
    d = Deck()
    d.title_slide("BÁO CÁO DỰ ÁN", "Biz Anchor Lending", "SPTD · 2026")
    d.section_slide("Phần 1", "Tổng quan tiến độ")
    s = d.content_slide("Tình hình triển khai", kicker="Cập nhật tuần")
    d.bullets(s, ["Điểm 1", "Điểm 2", ("Nhấn: ", "phần in đậm")])
    d.kpi_slide("Chỉ số chính", [("12", "Hồ sơ"), ("3", "Chờ duyệt"), ("95%", "Đúng hạn")])
    d.table_slide("Kế hoạch", [["Hạng mục","Chủ trì","Hạn"],["A","PO","30/6"]])
    d.closing_slide("Trân trọng cảm ơn")
    d.save("bao-cao.pptx")

Font: Roboto (chính). Nếu máy chưa cài, PowerPoint tự thay — cài Roboto để đúng chuẩn.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")

# ============ BRAND TOKENS (TPBank) ============
class C:
    PURPLE      = RGBColor(0x5E, 0x2E, 0x86)   # tím chủ đạo
    PURPLE_DARK = RGBColor(0x2C, 0x1E, 0x61)   # tím đậm / navy
    PURPLE_MID  = RGBColor(0x80, 0x77, 0xB2)   # tím trung
    PURPLE_SOFT = RGBColor(0x9D, 0x99, 0xB6)   # tím xám nhạt
    LAVENDER    = RGBColor(0xF1, 0xE9, 0xF7)   # tím rất nhạt (nền bảng / panel)
    LAVENDER2   = RGBColor(0xF3, 0xEC, 0xF8)
    GOLD        = RGBColor(0xFF, 0xC0, 0x00)   # gold accent
    ORANGE      = RGBColor(0xF7, 0x94, 0x1D)   # cam
    MAGENTA     = RGBColor(0xC0, 0x00, 0xC0)   # hồng/magenta accent
    INK         = RGBColor(0x21, 0x21, 0x21)   # text body
    RED         = RGBColor(0xFF, 0x00, 0x00)   # nhấn/cảnh báo
    GREEN       = RGBColor(0x00, 0xA0, 0x00)   # tích cực
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    CREAM       = RGBColor(0xFD, 0xFC, 0xF6)   # nền content sáng ấm
    GREY        = RGBColor(0x8A, 0x93, 0x9E)
    BORDER      = RGBColor(0xDD, 0xD6, 0xE6)

FONT = "Roboto"          # font chuẩn; fallback do PowerPoint xử lý
FONT_LIGHT = "Roboto"    # dùng Roboto cho cả tiêu đề (đồng bộ)

# Vị trí / kích thước chuẩn (cm → inch) lấy từ template gốc
def cm(v): return Inches(v / 2.54)
LOGO_W, LOGO_H = cm(4.97), cm(1.11)
LOGO_X, LOGO_Y = cm(1.16), cm(0.58)
ACCENT_Y = cm(1.95)      # dải accent ngay dưới header
MARGIN_X = cm(1.16)


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width  = cm(33.87)
        self.prs.slide_height = cm(19.05)
        self._blank = self.prs.slide_layouts[6]

    # ---------- primitives ----------
    def _slide(self):
        return self.prs.slides.add_slide(self._blank)

    def _rect(self, s, x, y, w, h, fill, line=None):
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None: sp.line.fill.background()
        else: sp.line.color.rgb = line; sp.line.width = Pt(0.75)
        sp.shadow.inherit = False
        return sp

    def _text(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, sa=6, ls=1.05):
        tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
        tf.word_wrap = True; tf.vertical_anchor = anchor
        for i, para in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0)
            p.line_spacing = ls
            for (t, sz, col, bold) in para:
                r = p.add_run(); r.text = t
                r.font.size = Pt(sz); r.font.color.rgb = col
                r.font.bold = bold; r.font.name = FONT
        return tb

    def _logo(self, s):
        s.shapes.add_picture(os.path.join(ASSETS, "logo-color.png"),
                             LOGO_X, LOGO_Y, LOGO_W, LOGO_H)

    def _accent(self, s, y=ACCENT_Y):
        # dải gradient gold→magenta thật (asset), full width
        s.shapes.add_picture(os.path.join(ASSETS, "accent-bar.png"),
                             0, y, self.prs.slide_width, cm(0.5))

    def _footer_bar(self, s):
        # thanh gold mảnh + khối chevron tím ở đáy (gợi motif brand)
        H = self.prs.slide_height; W = self.prs.slide_width
        self._rect(s, 0, H - cm(0.42), W, cm(0.42), C.GOLD)
        self._rect(s, W*0.62, H - cm(0.42), W*0.38, cm(0.42), C.PURPLE)

    # ---------- LAYOUT: TITLE ----------
    def title_slide(self, title, subtitle="", meta=""):
        s = self._slide()
        s.shapes.add_picture(os.path.join(ASSETS, "bg-title.jpg"),
                             0, 0, self.prs.slide_width, self.prs.slide_height)
        self._text(s, MARGIN_X, cm(5.5), cm(22), cm(5),
            [[(title, 40, C.WHITE, True)]], ls=1.02)
        if subtitle:
            self._text(s, MARGIN_X, cm(9.3), cm(22), cm(2),
                [[(subtitle, 24, C.GOLD, False)]])
        if meta:
            self._text(s, MARGIN_X, cm(11.2), cm(22), cm(1.5),
                [[(meta, 14, C.WHITE, False)]])
        return s

    # ---------- LAYOUT: SECTION DIVIDER ----------
    def section_slide(self, label, title):
        s = self._slide()
        s.shapes.add_picture(os.path.join(ASSETS, "bg-title.jpg"),
                             0, 0, self.prs.slide_width, self.prs.slide_height)
        self._text(s, MARGIN_X, cm(6.4), cm(22), cm(1.2),
            [[(label.upper(), 18, C.GOLD, True)]])
        self._text(s, MARGIN_X, cm(7.6), cm(24), cm(4),
            [[(title, 34, C.WHITE, True)]], ls=1.02)
        return s

    # ---------- LAYOUT: CONTENT (nền sáng) ----------
    def content_slide(self, title, kicker=""):
        """Trả về slide; dùng self.bullets / self.add_table / body_box để đổ nội dung."""
        s = self._slide()
        self._rect(s, 0, 0, self.prs.slide_width, self.prs.slide_height, C.CREAM)
        self._logo(s)
        if kicker:
            self._text(s, cm(16.5), LOGO_Y, cm(16), cm(1),
                [[(kicker.upper(), 11, C.GREY, True)]], align=PP_ALIGN.RIGHT,
                anchor=MSO_ANCHOR.MIDDLE)
        self._accent(s)
        self._text(s, MARGIN_X, cm(2.7), cm(30), cm(1.5),
            [[(title, 24, C.PURPLE, True)]])
        self._footer_bar(s)
        return s

    def body_box(self):
        """Vùng nội dung khả dụng (x, y, w, h) trên content slide."""
        return (MARGIN_X, cm(4.4), cm(31.5), cm(13.5))

    # ---------- content helpers ----------
    def bullets(self, s, items, x=None, y=None, w=None, size=16, gap=9,
                color=None):
        bx, by, bw, bh = self.body_box()
        x = bx if x is None else x; y = by if y is None else y
        w = bw if w is None else w; color = color or C.INK
        runs = []
        for it in items:
            if isinstance(it, tuple):
                lead, rest = it
                runs.append([("▸  ", size, C.PURPLE, True),
                             (lead, size, C.PURPLE, True), (rest, size, color, False)])
            else:
                runs.append([("▸  ", size, C.PURPLE, True), (it, size, color, False)])
        return self._text(s, x, y, w, bh, runs, sa=gap, ls=1.1)

    def panel(self, s, x, y, w, h, fill=None):
        return self._rect(s, x, y, w, h, fill or C.LAVENDER)

    # ---------- LAYOUT: KPI ----------
    def kpi_slide(self, title, kpis, kicker=""):
        s = self.content_slide(title, kicker)
        bx, by, bw, bh = self.body_box()
        n = len(kpis); gap = cm(0.5)
        cw = (bw - gap*(n-1)) / n
        ch = cm(4.2); y = by + cm(1.5)
        for i, (big, small) in enumerate(kpis):
            x = bx + i*(cw+gap)
            self._rect(s, x, y, cw, ch, C.LAVENDER)
            self._rect(s, x, y, cw, cm(0.16), C.GOLD)
            self._text(s, x, y+cm(0.6), cw, cm(2),
                [[(str(big), 40, C.PURPLE, True)]], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
            self._text(s, x, y+ch-cm(1.3), cw, cm(1.1),
                [[(small, 13, C.INK, False)]], align=PP_ALIGN.CENTER)
        return s

    # ---------- LAYOUT: TABLE ----------
    def table_slide(self, title, rows, col_widths=None, kicker=""):
        """rows[0] = header. rows: list[list[str]]."""
        s = self.content_slide(title, kicker)
        bx, by, bw, bh = self.body_box()
        nrows, ncols = len(rows), len(rows[0])
        from pptx.util import Cm
        gtbl = s.shapes.add_table(nrows, ncols, bx, by+cm(0.3), bw, cm(1.0)*nrows)
        tbl = gtbl.table
        if col_widths:
            total = sum(col_widths)
            for j, cwrel in enumerate(col_widths):
                tbl.columns[j].width = int(bw * cwrel / total)
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.margin_left = cm(0.2); cell.margin_right = cm(0.2)
                cell.margin_top = cm(0.08); cell.margin_bottom = cm(0.08)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                if i == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = C.PURPLE
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C.WHITE if i % 2 else C.LAVENDER
                tf = cell.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
                r = p.add_run(); r.text = str(val)
                r.font.name = FONT; r.font.size = Pt(12)
                r.font.bold = (i == 0)
                r.font.color.rgb = C.WHITE if i == 0 else C.INK
        return s

    # ---------- LAYOUT: CLOSING ----------
    def closing_slide(self, message="Trân trọng cảm ơn", sub=""):
        s = self._slide()
        s.shapes.add_picture(os.path.join(ASSETS, "bg-title.jpg"),
                             0, 0, self.prs.slide_width, self.prs.slide_height)
        self._text(s, MARGIN_X, cm(7.0), cm(28), cm(3),
            [[(message, 36, C.WHITE, True)]])
        if sub:
            self._text(s, MARGIN_X, cm(10.5), cm(28), cm(2),
                [[(sub, 18, C.GOLD, False)]])
        return s

    def save(self, path):
        self.prs.save(path)
        return path
